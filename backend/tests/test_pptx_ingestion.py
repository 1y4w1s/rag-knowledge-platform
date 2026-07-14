"""F2：PPT .pptx 入库管道 & prose chunk 测试。"""

from __future__ import annotations

import uuid
from pathlib import Path

import pytest
from httpx import AsyncClient
from sqlalchemy import select

from app.core.database import SessionLocal
from app.models.document import Document
from app.models.document_chunk import DocumentChunk
from app.models.enums import DocumentStatus
from app.services.ingestion.pipeline import process_document_ingestion

from tests.conftest import create_test_kb as _create_kb

pytestmark = pytest.mark.asyncio


def _make_golden_pptx(path: Path) -> None:
    """Create a small test pptx with 3 slides."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
    txBox.text_frame.text = "知岸产品介绍"
    txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
    txBox2.text_frame.text = "企业级知识库 RAG 平台"
    # Add notes
    notes_slide = slide1.notes_slide
    notes_slide.notes_text_frame.text = "开场介绍知岸的核心定位"

    # Slide 2: Features
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox3 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
    txBox3.text_frame.text = "核心功能"
    txBox4 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(3))
    txBox4.text_frame.text = "多格式文档上传\n智能切片与检索\n溯源对话引用"

    # Slide 3: Q&A
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox5 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
    txBox5.text_frame.text = "Q&A"
    txBox6 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
    txBox6.text_frame.text = "欢迎提问"

    prs.save(str(path))


async def _load_chunks(document_id: uuid.UUID) -> list[DocumentChunk]:
    async with SessionLocal() as db:
        result = await db.scalars(
            select(DocumentChunk)
            .where(DocumentChunk.document_id == document_id)
            .order_by(DocumentChunk.chunk_index)
        )
        return list(result.all())


@pytest.fixture
def upload_dir(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    from app.core.config import settings
    monkeypatch.setattr(settings, "upload_dir", str(tmp_path))
    return tmp_path


async def _ingest_fixture(
    *,
    kb_id: uuid.UUID,
    user_id: uuid.UUID,
    source: Path,
    file_type: str,
    upload_dir: Path,
) -> Document:
    doc_id = uuid.uuid4()
    storage_dir = upload_dir / str(kb_id) / str(doc_id)
    storage_dir.mkdir(parents=True, exist_ok=True)
    storage_path = storage_dir / f"{uuid.uuid4()}.{file_type}"
    storage_path.write_bytes(source.read_bytes())

    async with SessionLocal() as db:
        doc = Document(
            id=doc_id,
            kb_id=kb_id,
            filename=source.name,
            file_type=file_type,
            file_size=storage_path.stat().st_size,
            storage_path=str(storage_path),
            status=DocumentStatus.queued,
            uploaded_by=user_id,
        )
        db.add(doc)
        await db.commit()

    await process_document_ingestion(doc_id)

    async with SessionLocal() as db:
        row = await db.get(Document, doc_id)
        assert row is not None
        return row


async def test_pptx_ingestion_creates_prose_chunks(
    client: AsyncClient,
    register_and_login,
    upload_dir: Path,
    tmp_path: Path,
) -> None:
    """pptx 文件入库后生成 prose chunk，每 slide 一段。"""
    pytest.importorskip("pptx")
    pptx_path = tmp_path / "golden_deck.pptx"
    _make_golden_pptx(pptx_path)

    headers, user = await register_and_login(prefix="pptx-test")
    kb = await _create_kb(client, headers, user)
    kb_id = uuid.UUID(kb["id"])
    user_id = uuid.UUID(user["id"])

    doc = await _ingest_fixture(
        kb_id=kb_id,
        user_id=user_id,
        source=pptx_path,
        file_type="pptx",
        upload_dir=upload_dir,
    )

    assert doc.status == DocumentStatus.completed
    assert doc.chunk_count is not None and doc.chunk_count > 0
    assert doc.processing_completed_at is not None

    chunks = await _load_chunks(doc.id)
    assert len(chunks) == doc.chunk_count
    assert all(c.embedding is not None for c in chunks)

    # Verify slide titles appear in heading_path (stored as section_title)
    headings = " ".join(c.heading_path for c in chunks if c.heading_path)
    assert "知岸产品介绍" in headings
    assert "核心功能" in headings
    assert "Q&A" in headings

    # Verify notes were included in content
    contents = " ".join(c.content for c in chunks)
    assert "开场介绍" in contents

    # Verify page_number (slide index) is recorded
    chunks_with_page = [c for c in chunks if c.page_number is not None]
    assert len(chunks_with_page) > 0


async def test_pptx_ingestion_upload_endpoint(
    client: AsyncClient,
    register_and_login,
    tmp_path: Path,
) -> None:
    """通过上传 API 验证 pptx 可正常上传入库。"""
    pytest.importorskip("pptx")
    pptx_path = tmp_path / "test_upload.pptx"
    _make_golden_pptx(pptx_path)

    headers, user = await register_and_login(prefix="pptx-upload")
    kb = await _create_kb(client, headers, user)
    kb_id = kb["id"]

    # Upload via API
    with open(pptx_path, "rb") as f:
        upload_resp = await client.post(
            f"/api/v1/knowledge-bases/{kb_id}/documents",
            files={"files": ("test_deck.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            headers=headers,
        )
    assert upload_resp.status_code == 201, upload_resp.text
    data = upload_resp.json()
    docs = data["documents"]
    assert len(docs) == 1
    assert docs[0]["file_type"] == "pptx"
    assert docs[0]["status"] in ("queued", "completed")
