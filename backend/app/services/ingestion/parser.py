"""多格式文档解析（TECH-4.2）。"""

from __future__ import annotations

import re
from pathlib import Path

from app.services.ingestion.parser_pdf import (
    _merge_cross_page_blocks,
    detect_scanned_pdf,
    parse_pdf,
    parse_pdf_ocr,
)
from app.services.ingestion.table_detection import is_markdown_table_block
from app.services.ingestion.types import ParsedBlock

CHAPTER_RE = re.compile(
    r"^(?:第[一二三四五六七八九十百千\d]+章[^\n]{0,40}|"
    r"Chapter\s+\d+[^\n]{0,40}|"
    r"\d+(?:\.\d+)*\s+[^\n]{1,60}|"
    r"#{1,3}\s+.+)$"
)
MD_HEADER_RE = re.compile(r"^(#{1,3})\s+(.+)$")


def parse_txt(content: str) -> list[ParsedBlock]:
    blocks: list[ParsedBlock] = []
    heading_stack: list[str] = []
    doc_title: str | None = None

    for para in re.split(r"\n\s*\n", content.strip()):
        text = para.strip()
        if not text:
            continue

        if CHAPTER_RE.match(text.split("\n", 1)[0]):
            title_line = text.split("\n", 1)[0].strip()
            title = re.sub(r"^#+\s*", "", title_line)
            if doc_title is None and not title.startswith("第") and "." not in title[:4]:
                doc_title = title
                heading_stack = [title]
            else:
                heading_stack = ([doc_title] if doc_title else []) + [title]
            body = text.split("\n", 1)[1].strip() if "\n" in text else ""
            if body:
                blocks.append(
                    ParsedBlock(
                        content=body,
                        section_title=title,
                        heading_path=">".join(heading_stack),
                    )
                )
            continue

        blocks.append(
            ParsedBlock(
                content=text,
                section_title=heading_stack[-1] if heading_stack else None,
                heading_path=">".join(heading_stack) if heading_stack else None,
                block_kind="table" if is_markdown_table_block(text) else "prose",
            )
        )
    return blocks


def parse_md(content: str) -> list[ParsedBlock]:
    blocks: list[ParsedBlock] = []
    heading_stack: list[tuple[int, str]] = []
    doc_title: str | None = None

    for para in re.split(r"\n\s*\n", content.strip()):
        text = para.strip()
        if not text:
            continue

        header_match = MD_HEADER_RE.match(text.split("\n", 1)[0])
        if header_match:
            level = len(header_match.group(1))
            title = header_match.group(2).strip()
            if level == 1 and doc_title is None:
                doc_title = title
            heading_stack = [(lvl, name) for lvl, name in heading_stack if lvl < level]
            heading_stack.append((level, title))
            path = ">".join(name for _, name in heading_stack)
            body = text.split("\n", 1)[1].strip() if "\n" in text else ""
            if body:
                blocks.append(
                    ParsedBlock(
                        content=body,
                        section_title=title,
                        heading_path=path,
                    )
                )
            continue

        path = ">".join(name for _, name in heading_stack) if heading_stack else None
        kind = "table" if is_markdown_table_block(text) else "prose"
        blocks.append(
            ParsedBlock(
                content=text,
                section_title=heading_stack[-1][1] if heading_stack else None,
                heading_path=path,
                block_kind=kind,
            )
        )
    return blocks


def _docx_table_to_markdown(table) -> str:
    rows: list[str] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    return "\n".join(rows)


def _iter_docx_body(doc) -> list:
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in doc.element.body:
        if isinstance(child, CT_P):
            yield Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield Table(child, doc)


def parse_docx(path: Path) -> list[ParsedBlock]:
    from docx import Document as DocxDocument

    blocks: list[ParsedBlock] = []
    heading_stack: list[str] = []
    doc_title: str | None = None

    doc = DocxDocument(str(path))
    buffer: list[str] = []
    current_meta = ParsedBlock(content="")

    def flush() -> None:
        nonlocal buffer, current_meta
        if not buffer:
            return
        text = "\n".join(buffer).strip()
        if text:
            blocks.append(
                ParsedBlock(
                    content=text,
                    section_title=current_meta.section_title,
                    heading_path=current_meta.heading_path,
                    block_kind="prose",
                )
            )
        buffer = []

    for item in _iter_docx_body(doc):
        if hasattr(item, "rows"):
            flush()
            table_text = _docx_table_to_markdown(item)
            if table_text.strip():
                blocks.append(
                    ParsedBlock(
                        content=table_text,
                        section_title=current_meta.section_title,
                        heading_path=current_meta.heading_path,
                        block_kind="table",
                    )
                )
            continue

        para = item
        text = para.text.strip()
        if not text:
            flush()
            continue

        style_name = para.style.name if para.style else ""
        if style_name.startswith("Heading"):
            flush()
            try:
                level = int(style_name.replace("Heading", "").strip() or "1")
            except ValueError:
                level = 1
            if level == 1 and doc_title is None:
                doc_title = text
            heading_stack = heading_stack[: max(0, level - 1)]
            heading_stack.append(text)
            current_meta = ParsedBlock(
                content="",
                section_title=text,
                heading_path=">".join(heading_stack),
            )
            continue

        if CHAPTER_RE.match(text):
            flush()
            if doc_title is None:
                doc_title = text
            heading_stack = [doc_title] if doc_title else []
            if text != doc_title:
                heading_stack.append(text)
            current_meta = ParsedBlock(
                content="",
                section_title=text,
                heading_path=">".join(heading_stack),
            )
            continue

        buffer.append(text)

    flush()
    return blocks


def _rows_to_markdown_table(rows: list[tuple]) -> list[str]:
    """Convert a 2D array of values to pipe-markdown table lines."""
    if not rows:
        return []
    lines: list[str] = []
    # Header row
    lines.append("| " + " | ".join(str(c or "") for c in rows[0]) + " |")
    # Separator
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    # Data rows
    for row in rows[1:]:
        lines.append("| " + " | ".join(str(c or "") for c in row) + " |")
    return lines


def parse_xlsx(path: Path) -> list[ParsedBlock]:
    """Excel: each sheet as a markdown table block."""
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    blocks: list[ParsedBlock] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows or all(cell is None for cell in rows[0]):
            continue
        md_lines = _rows_to_markdown_table(rows)
        blocks.append(
            ParsedBlock(
                content="\n".join(md_lines),
                section_title=sheet_name,
                heading_path=sheet_name,
                block_kind="table",
            )
        )
    wb.close()
    return blocks


def parse_pptx(path: Path) -> list[ParsedBlock]:
    """PPT: each slide as a prose block, notes appended to body."""
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list[ParsedBlock] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if title_shape is not None and shape == title_shape:
                title = text
            else:
                body_parts.append(text)
        # Fallback: first textbox as title when blank layout (no title placeholder)
        if not title and body_parts:
            title = body_parts.pop(0)
        # Notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                body_parts.append(f"【备注】{notes}")
        content = "\n\n".join(body_parts) if body_parts else ""
        heading_path = title or f"Slide {i}"
        blocks.append(
            ParsedBlock(
                content=content or heading_path,
                section_title=title or None,
                heading_path=heading_path,
                page_number=i,
                block_kind="prose",
            )
        )
    return blocks


def parse_document(path: Path, file_type: str, *, pdf_batch_pages: int = 10) -> list[ParsedBlock]:
    ext = file_type.lower().lstrip(".")

    # 魔数检测：文件扩展名与实际内容是否匹配
    _MAGIC_SIGNATURES = {
        "pdf": b"%PDF",
        "docx": b"PK",
        "xlsx": b"PK",
        "pptx": b"PK",
    }
    expected_magic = _MAGIC_SIGNATURES.get(ext)
    if expected_magic is not None:
        header = path.read_bytes()[:4]
        if not header.startswith(expected_magic):
            raise ValueError(f"文件格式不匹配：扩展名为 .{ext} 但内容格式不符")

    if ext == "pdf":
        if detect_scanned_pdf(path):
            from app.services.ingestion.ocr import is_ocr_enabled, is_ocr_runtime_available

            if not is_ocr_enabled():
                raise ValueError("不支持扫描件")
            if not is_ocr_runtime_available():
                raise ValueError("OCR 服务未启用")
            return parse_pdf_ocr(path)

        blocks = parse_pdf(path, batch_pages=pdf_batch_pages)
        if not blocks:
            # 文字层为空时尝试 OCR fallback
            from app.services.ingestion.ocr import is_ocr_enabled, is_ocr_runtime_available
            if is_ocr_enabled() and is_ocr_runtime_available():
                blocks = parse_pdf_ocr(path)
            if not blocks:
                raise ValueError("不支持扫描件")
        return blocks
    if ext == "txt":
        return parse_txt(path.read_text(encoding="utf-8"))
    if ext == "md":
        return parse_md(path.read_text(encoding="utf-8"))
    if ext == "docx":
        return parse_docx(path)
    if ext == "xlsx":
        return parse_xlsx(path)
    if ext == "pptx":
        return parse_pptx(path)
    if ext in ("png", "jpg", "jpeg"):
        return parse_image_ocr(path)
    raise ValueError(f"不支持的文件类型: {file_type}")


IMAGE_EXTENSIONS = {"png", "jpg", "jpeg"}


def parse_image_ocr(path: Path) -> list[ParsedBlock]:
    """图片文件走 OCR 识别。"""
    from app.services.ingestion.ocr import is_ocr_enabled, is_ocr_runtime_available, ocr_image_path

    if not is_ocr_enabled():
        raise ValueError("OCR 未启用")
    if not is_ocr_runtime_available():
        raise ValueError("OCR 服务未安装")

    page_number, text = ocr_image_path(path)
    cleaned = text.strip()
    if not cleaned:
        raise ValueError("OCR 未识别到文字")
    return [ParsedBlock(content=cleaned, page_number=page_number)]
