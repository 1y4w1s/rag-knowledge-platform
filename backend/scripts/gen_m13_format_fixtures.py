#!/usr/bin/env python3
"""Generate M13 format-acceptance fixtures (NW-36).

Run from repo root or backend/:

    python backend/scripts/gen_m13_format_fixtures.py
    python scripts/gen_m13_format_fixtures.py

Writes to backend/tests/fixtures/m13_format/. Safe to re-run (overwrite).
"""

from __future__ import annotations

from pathlib import Path

_ROOT = Path(__file__).resolve().parents[1]
_OUT = _ROOT / "tests" / "fixtures" / "m13_format"

# Unique needles for citation smoke (keep in sync with eval-M13-format-matrix.md)
NEEDLES = {
    "txt": "M13-NEEDLE-TXT-8841",
    "md": "M13-NEEDLE-MD-8842",
    "docx": "M13-NEEDLE-DOCX-8843",
    "pdf": "M13-NEEDLE-PDF-8844",
    "xlsx": "M13-NEEDLE-XLSX-8845",
    "pptx": "M13-NEEDLE-PPTX-8846",
}


def _write_txt(path: Path) -> None:
    path.write_text(
        "Ruige M13 format matrix sample (plain text).\n"
        f"Acceptance code: {NEEDLES['txt']}\n"
        "This line exists so chat can cite the needle.\n",
        encoding="utf-8",
    )


def _write_md(path: Path) -> None:
    path.write_text(
        "# M13 Markdown sample\n\n"
        f"Acceptance code: **{NEEDLES['md']}**\n\n"
        "Use this file to verify markdown ingestion and citations.\n",
        encoding="utf-8",
    )


def _write_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("M13 DOCX sample", level=1)
    doc.add_paragraph(
        f"Acceptance code: {NEEDLES['docx']}. "
        "Cite this paragraph in chat smoke."
    )
    doc.save(str(path))


def _write_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    c.drawString(72, 720, "Ruige M13 PDF text-layer sample")
    c.drawString(72, 700, f"Acceptance code: {NEEDLES['pdf']}")
    c.drawString(72, 680, "Enough extractable text for non-OCR PDF path.")
    c.showPage()
    c.save()


def _write_xlsx(path: Path) -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "M13"
    ws.append(["item", "code", "note"])
    ws.append(["format_matrix", NEEDLES["xlsx"], "xlsx citation needle"])
    ws.append(["owner", "ops", "NW-36"])
    wb.save(str(path))


def _write_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(2))
    box.text_frame.text = f"M13 PPTX sample\nAcceptance code: {NEEDLES['pptx']}"
    prs.save(str(path))


def _write_scan_readme(path: Path) -> None:
    path.write_text(
        "# Optional scanned PDF (OCR)\n\n"
        "Do **not** commit a blank scanned PDF as a gate fixture.\n"
        "For optional smoke: create a nearly blank PDF (text layer empty),\n"
        "enable OCR runtime (M11-B2), upload, ask a question about any OCR text.\n"
        "See docs/tasks/eval-M13-format-matrix.md § optional.\n",
        encoding="utf-8",
    )


def main() -> None:
    _OUT.mkdir(parents=True, exist_ok=True)
    _write_txt(_OUT / "m13_plain.txt")
    _write_md(_OUT / "m13_notes.md")
    _write_docx(_OUT / "m13_handbook.docx")
    _write_pdf(_OUT / "m13_text_layer.pdf")
    _write_xlsx(_OUT / "m13_ledger.xlsx")
    _write_pptx(_OUT / "m13_deck.pptx")
    _write_scan_readme(_OUT / "README_SCAN_OPTIONAL.md")
    print(f"Wrote fixtures under {_OUT}")
    for k, v in NEEDLES.items():
        print(f"  {k}: {v}")


if __name__ == "__main__":
    main()
